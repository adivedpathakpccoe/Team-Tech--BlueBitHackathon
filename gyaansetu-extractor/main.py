from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import os
import shutil
import tempfile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl

app = FastAPI(title="GYAANSETU Content Extractor")

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_pdf(file_path):
    text = []
    with fitz.open(file_path) as doc:
        for i, page in enumerate(doc):
            page_text = page.get_text().strip()
            if page_text:
                text.append(f"--- Page {i + 1} ---\n{page_text}")
    return "\n\n".join(text)

def extract_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        combined_slide = "\n".join(slide_text).strip()
        if combined_slide:
            text.append(f"--- Slide {i + 1} ---\n{combined_slide}")
            
    return "\n\n".join(text)

def extract_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"--- Sheet: {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)

@app.post("/extract")
async def extract_content(file: UploadFile = File(...)):
    suffix = os.path.splitext(file.filename)[1].lower()
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        try:
            shutil.copyfileobj(file.file, tmp)
            tmp_path = tmp.name
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to save uploaded file: {str(e)}")

    try:
        content = ""
        if suffix == ".pdf":
            content = extract_pdf(tmp_path)
        elif suffix == ".docx":
            content = extract_docx(tmp_path)
        elif suffix == ".pptx":
            content = extract_pptx(tmp_path)
        elif suffix in [".xlsx", ".xls"]:
            content = extract_xlsx(tmp_path)
        elif suffix in [".txt", ".md", ".csv"]:
            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        else:
            return {
                "filename": file.filename,
                "error": f"Unsupported file type: {suffix}",
                "success": False
            }

        return {
            "filename": file.filename,
            "content": content.strip(),
            "success": True
        }
    except Exception as e:
        return {
            "filename": file.filename,
            "error": f"Extraction failed: {str(e)}",
            "success": False
        }
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

@app.get("/health")
async def health_check():
    return {"status": "healthy"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8001, reload=True)
